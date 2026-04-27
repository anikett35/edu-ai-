"""
services/file_service.py
────────────────────────
Handles the full material-upload pipeline:

  1. Validate file type & size
  2. Save to local filesystem (uploads/<teacher_id>/<filename>)
  3. Extract text — routed by file type:
       .pdf          → PyMuPDF (fitz)
       .pptx / .ppt  → python-pptx  (slide text + speaker notes)
       .docx / .doc  → python-docx  (paragraphs + tables)
  4. Clean extracted text
  5. Split into chunks (200-500 words each)
  6. Generate sentence-transformer embeddings for each chunk
  7. Return structured data ready for MongoDB insertion

Supported formats:
  PDF  (.pdf)  — lecture notes, research papers, textbooks
  PPTX (.pptx, .ppt) — slide decks
  DOCX (.docx, .doc) — Word documents, study guides

Why chunking matters:
  Vector similarity search works best on focused, coherent passages.
  A 50-page PDF embedded as one vector loses granular context.
  Chunked retrieval lets the RAG system pinpoint exactly which passage
  answers a student's question.

Why sentence-transformers/all-MiniLM-L6-v2?
  - 384-dimension embeddings (compact, fast)
  - Trained on 1B+ sentence pairs
  - Excellent semantic similarity for Q&A retrieval
  - Runs fully locally — no API calls

Extra dependencies needed (add to requirements.txt):
  python-pptx==0.6.23
  python-docx==1.1.2
"""

import os
import re
import logging
from pathlib import Path
from typing import List, Tuple

import fitz  # PyMuPDF — PDF extraction
import numpy as np
from sentence_transformers import SentenceTransformer
from fastapi import UploadFile, HTTPException, status
from pptx import Presentation                          # python-pptx — PPTX extraction
from pptx.util import Pt
import docx                                            # python-docx — DOCX extraction

from core.config import settings

logger = logging.getLogger(__name__)

# ── Embedding model (singleton – loaded once at import) ──────────────────────
# Loading is slow (~2s); we do it at startup not per request.
_embedding_model: SentenceTransformer | None = None


def get_embedding_model() -> SentenceTransformer:
    global _embedding_model
    if _embedding_model is None:
        logger.info("Loading embedding model: %s", settings.embedding_model)
        _embedding_model = SentenceTransformer(settings.embedding_model)
        logger.info("✅  Embedding model loaded")
    return _embedding_model


# ── Constants ────────────────────────────────────────────────────────────────

ALLOWED_EXTENSIONS = {".pdf", ".pptx", ".ppt", ".docx", ".doc"}
# MIME types mapped to extensions for double-checking content-type header
ALLOWED_MIME_TYPES = {
    "application/pdf",
    "application/vnd.openxmlformats-officedocument.presentationml.presentation",  # .pptx
    "application/vnd.ms-powerpoint",                                               # .ppt
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document",     # .docx
    "application/msword",                                                          # .doc
}
MAX_BYTES = settings.max_file_size_mb * 1024 * 1024
CHUNK_SIZE_WORDS = 350                 # target words per chunk
CHUNK_OVERLAP_WORDS = 50               # overlap to preserve context at boundaries


# ════════════════════════════════════════════════════════════════════════════
# PUBLIC API
# ════════════════════════════════════════════════════════════════════════════

async def process_uploaded_file(
    file: UploadFile,
    teacher_id: str,
    subject: str,
) -> dict:
    """
    Full pipeline: validate → save → extract → chunk → embed.
    Extraction is routed by file extension:
      .pdf         → PyMuPDF
      .pptx/.ppt   → python-pptx
      .docx/.doc   → python-docx
    Returns dict ready to insert into MongoDB materials collection.
    """
    # Step 1 – validate
    _validate_file(file)

    # Step 2 – save
    file_path = await _save_file(file, teacher_id)

    # Step 3 – extract (route by extension)
    suffix = file_path.suffix.lower()
    try:
        if suffix == ".pdf":
            raw_text = _extract_text_from_pdf(file_path)
        elif suffix in {".pptx", ".ppt"}:
            raw_text = _extract_text_from_pptx(file_path)
        elif suffix in {".docx", ".doc"}:
            raw_text = _extract_text_from_docx(file_path)
        else:
            # Should never reach here after _validate_file, but defensive
            raise HTTPException(
                status_code=status.HTTP_415_UNSUPPORTED_MEDIA_TYPE,
                detail=f"Unsupported file type: {suffix}",
            )
    except HTTPException:
        raise
    except Exception as exc:
        logger.error("Text extraction failed for %s: %s", file_path, exc)
        raise HTTPException(
            status_code=status.HTTP_422_UNPROCESSABLE_ENTITY,
            detail=f"Failed to extract text from file: {exc}",
        )

    if not raw_text.strip():
        raise HTTPException(
            status_code=status.HTTP_422_UNPROCESSABLE_ENTITY,
            detail=(
                f"Could not extract any text from the uploaded {suffix.upper()} file. "
                "For PDFs: ensure it is not a scanned image-only document. "
                "For PPTX: ensure slides contain text boxes, not just images. "
                "For DOCX: ensure the document contains readable paragraph text."
            ),
        )

    # Step 4 – clean
    cleaned_text = _clean_text(raw_text)

    # Step 5 – chunk
    chunks = _split_into_chunks(cleaned_text, CHUNK_SIZE_WORDS, CHUNK_OVERLAP_WORDS)
    logger.info("📄  %d chunks created from '%s'", len(chunks), file.filename)

    # Step 6 – embed
    embeddings = _embed_chunks(chunks)

    return {
        "teacher_id": teacher_id,
        "subject": subject,
        "file_name": file.filename,
        "file_path": str(file_path),
        "file_type": suffix.lstrip("."),       # "pdf" | "pptx" | "docx" etc.
        "extracted_chunks": chunks,
        "embeddings": embeddings,
        "chunk_count": len(chunks),
    }


def embed_query(question: str) -> np.ndarray:
    """Embed a single question string for RAG similarity search."""
    model = get_embedding_model()
    return model.encode(question, normalize_embeddings=True)


# ════════════════════════════════════════════════════════════════════════════
# INTERNAL HELPERS
# ════════════════════════════════════════════════════════════════════════════

def _validate_file(file: UploadFile) -> None:
    """
    Raise HTTPException if:
      - File extension is not in ALLOWED_EXTENSIONS
      - File size exceeds MAX_BYTES (if Content-Length is present)

    Note: .ppt and .doc are legacy binary formats. python-pptx and python-docx
    both handle them via internal format detection, but extraction quality is
    lower than the XML-based .pptx/.docx formats. Teachers should be encouraged
    to use modern formats where possible.
    """
    if not file.filename:
        raise HTTPException(
            status_code=status.HTTP_422_UNPROCESSABLE_ENTITY,
            detail="File has no name",
        )
    suffix = Path(file.filename).suffix.lower()
    if suffix not in ALLOWED_EXTENSIONS:
        raise HTTPException(
            status_code=status.HTTP_415_UNSUPPORTED_MEDIA_TYPE,
            detail=(
                f"Unsupported file type '{suffix}'. "
                f"Allowed types: PDF (.pdf), PowerPoint (.pptx, .ppt), Word (.docx, .doc)"
            ),
        )
    if file.size and file.size > MAX_BYTES:
        raise HTTPException(
            status_code=status.HTTP_413_REQUEST_ENTITY_TOO_LARGE,
            detail=f"File exceeds {settings.max_file_size_mb} MB limit",
        )


async def _save_file(file: UploadFile, teacher_id: str) -> Path:
    """Save upload to uploads/<teacher_id>/<filename>. Returns absolute path."""
    dest_dir = Path(settings.upload_dir) / teacher_id
    dest_dir.mkdir(parents=True, exist_ok=True)

    # Sanitise filename
    safe_name = re.sub(r"[^\w\.\-]", "_", file.filename)
    file_path = dest_dir / safe_name

    content = await file.read()

    # Enforce size after reading full content
    if len(content) > MAX_BYTES:
        raise HTTPException(
            status_code=status.HTTP_413_REQUEST_ENTITY_TOO_LARGE,
            detail=f"File exceeds {settings.max_file_size_mb} MB limit",
        )

    with open(file_path, "wb") as f:
        f.write(content)

    logger.info("💾  File saved: %s (%d bytes)", file_path, len(content))
    return file_path


def _extract_text_from_pdf(file_path: Path) -> str:
    """
    PyMuPDF (fitz) — best-in-class PDF text extraction.
    Handles PDF/A, partially corrupted files, and multi-column layouts.
    Each page's text is separated by a newline.
    """
    doc = fitz.open(str(file_path))
    pages_text = []
    for page in doc:
        pages_text.append(page.get_text("text"))
    doc.close()
    return "\n".join(pages_text)


def _extract_text_from_pptx(file_path: Path) -> str:
    """
    python-pptx extractor — handles .pptx and .ppt files.

    Extraction strategy:
      1. Iterate every slide in order
      2. For each slide, collect all text from:
           a. Text frames inside shapes (titles, content boxes, text boxes)
           b. Table cells (common in academic slides)
           c. Speaker notes (often contain additional explanation — valuable for RAG)
      3. Prefix each slide's content with "Slide N:" so chunks retain context

    Why include speaker notes?
      Teachers often put detailed explanations in notes that don't appear
      on the slide itself. These are extremely valuable for RAG retrieval.
    """
    prs = Presentation(str(file_path))
    slide_texts: List[str] = []

    for slide_num, slide in enumerate(prs.slides, start=1):
        parts: List[str] = [f"[Slide {slide_num}]"]

        for shape in slide.shapes:
            # Regular text frames (titles, content, text boxes)
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = para.text.strip()
                    if line:
                        parts.append(line)

            # Tables inside slides
            if shape.has_table:
                for row in shape.table.rows:
                    row_text = " | ".join(
                        cell.text.strip() for cell in row.cells if cell.text.strip()
                    )
                    if row_text:
                        parts.append(row_text)

        # Speaker notes (separate notes_slide object)
        if slide.has_notes_slide:
            notes_frame = slide.notes_slide.notes_text_frame
            if notes_frame:
                notes_text = notes_frame.text.strip()
                if notes_text and notes_text.lower() != "click to add notes":
                    parts.append(f"[Notes: {notes_text}]")

        if len(parts) > 1:  # more than just the slide header
            slide_texts.append("\n".join(parts))

    full_text = "\n\n".join(slide_texts)
    logger.info("📊  PPTX: extracted %d slides from '%s'", len(prs.slides), file_path.name)
    return full_text


def _extract_text_from_docx(file_path: Path) -> str:
    """
    python-docx extractor — handles .docx and .doc files.

    Extraction strategy:
      1. All paragraphs in document body (preserving order)
      2. All tables: each row joined with pipe separator
      3. Headers are detected by style name and prefixed with "##" to
         signal section boundaries — useful if you later want
         header-aware chunking.

    Why extract tables?
      DOCX study materials frequently contain comparison tables,
      formula tables, and data summaries. Ignoring them loses content.
    """
    doc = docx.Document(str(file_path))
    parts: List[str] = []

    # _iter_block_items preserves paragraph/table interleave order
    # (doc.paragraphs and doc.tables lose ordering relative to each other)
    for block in _iter_block_items(doc):
        if isinstance(block, docx.text.paragraph.Paragraph):
            text = block.text.strip()
            if not text:
                continue
            # Detect headings by style
            style_name = block.style.name.lower() if block.style else ""
            if "heading" in style_name:
                parts.append(f"\n## {text}")
            else:
                parts.append(text)

        elif isinstance(block, docx.table.Table):
            for row in block.rows:
                row_text = " | ".join(
                    cell.text.strip() for cell in row.cells if cell.text.strip()
                )
                if row_text:
                    parts.append(row_text)

    full_text = "\n".join(parts)
    logger.info("📝  DOCX: extracted %d blocks from '%s'", len(parts), file_path.name)
    return full_text


def _iter_block_items(doc: "docx.Document"):
    """
    Yield paragraphs and tables from a DOCX document in document order.
    python-docx's .paragraphs and .tables properties return them separately,
    losing interleaving. This generator preserves the original order by
    walking the XML body directly.
    """
    from docx.oxml.ns import qn
    from docx.table import Table
    from docx.text.paragraph import Paragraph

    body = doc.element.body
    for child in body.iterchildren():
        if child.tag == qn("w:p"):
            yield Paragraph(child, doc)
        elif child.tag == qn("w:tbl"):
            yield Table(child, doc)


def _clean_text(text: str) -> str:
    """
    Remove noise introduced by PDF extraction:
      - Control characters
      - Repeated whitespace / newlines
      - Page headers / footers (heuristic: very short lines < 4 words)
      - Ligature artifacts (ﬁ → fi, etc.)
    """
    # Unicode ligatures
    replacements = {"ﬁ": "fi", "ﬂ": "fl", "ﬀ": "ff", "ﬃ": "ffi", "ﬄ": "ffl"}
    for src, dst in replacements.items():
        text = text.replace(src, dst)

    # Remove control characters except newlines
    text = re.sub(r"[^\x09\x0A\x0D\x20-\x7E\u00A0-\uFFFF]", " ", text)

    # Collapse multiple spaces
    text = re.sub(r"[ \t]+", " ", text)

    # Remove lines shorter than 4 words (likely headers/footers/page numbers)
    lines = text.splitlines()
    lines = [ln for ln in lines if len(ln.split()) >= 4]

    # Collapse excessive blank lines
    text = "\n".join(lines)
    text = re.sub(r"\n{3,}", "\n\n", text)

    return text.strip()


def _split_into_chunks(
    text: str, chunk_size: int, overlap: int
) -> List[str]:
    """
    Word-level sliding-window chunking.

    Why overlap?
      A sentence at the boundary of chunk N might be key to answering a question.
      Overlapping by `overlap` words ensures no sentence is cut off cold.
    """
    words = text.split()
    chunks: List[str] = []
    start = 0

    while start < len(words):
        end = min(start + chunk_size, len(words))
        chunk = " ".join(words[start:end])
        chunks.append(chunk)
        if end == len(words):
            break
        start += chunk_size - overlap  # slide forward with overlap

    return chunks


def _embed_chunks(chunks: List[str]) -> List[List[float]]:
    """
    Generate normalised embeddings for each chunk.
    Returns list of plain Python lists (JSON-serialisable for MongoDB storage).

    normalize_embeddings=True → vectors on unit sphere → cosine sim = dot product
    This halves computation cost during retrieval.
    """
    model = get_embedding_model()
    embeddings: np.ndarray = model.encode(
        chunks,
        batch_size=32,
        normalize_embeddings=True,
        show_progress_bar=False,
    )
    return embeddings.tolist()